import os
import win32com.client as win32
from pdf2docx import Converter
import fitz
import os
from pptx import Presentation
from pypdf import PdfReader, PdfWriter

def convertir_word_a_pdf(word_path: str, pdf_path: str) -> None:
    """
    Convierte un archivo Word a PDF usando la API de Microsoft Word.
    word_path: ruta absoluta al archivo .doc/.docx
    pdf_path:  ruta absoluta de salida .pdf
    """
    word_app = win32.Dispatch("Word.Application")
    word_app.Visible = False
    try:
        doc = word_app.Documents.Open(word_path)
        doc.SaveAs(pdf_path, FileFormat=17)  # 17 = wdFormatPDF
        doc.Close()
    finally:
        word_app.Quit()


def convertir_pdf_word(pdf,word):
    pdf = os.path.abspath(pdf)
    word = os.path.abspath(word)

    cv = Converter(pdf)
    cv.convert(word, start=0, end=None)
    cv.close()


def convertir_pptx_pdf(pptx,pdf):
    pptx = os.path.abspath(pptx)
    pdf = os.path.abspath(pdf)

    powerpoint = win32.Dispatch("PowerPoint.Application")
    powerpoint.Visible = False

    presentation = powerpoint.Presentations.Open(pptx)
    presentation.SaveAs(pdf, FileFormat=32)
    presentation.Close()



def convertir_pdf_a_pptx(pdf, pptx):
    import tempfile
    pdf = os.path.abspath(pdf)
    pptx = os.path.abspath(pptx)
    doc = fitz.open(pdf)
    prs = Presentation()
    with tempfile.TemporaryDirectory() as temp_dir:
        for i, pagina in enumerate(doc):
            pix = pagina.get_pixmap(matrix=fitz.Matrix(2, 2))
            imagen = os.path.join(temp_dir, f"page_{i}.png")
            pix.save(imagen)
            slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(slide_layout)
            slide.shapes.add_picture(imagen, 0, 0, width=prs.slide_width, height=prs.slide_height)
    prs.save(pptx)

def unir_pdf(pdf):
    writer = PdfWriter()

    for archivo in pdf:
        reader = PdfReader(archivo)
        for pagina in reader.pages:
            writer.add_page(pagina)

    return writer


def extraer_paginas(pdf, salida, inicio, fin):
    reader = PdfReader(pdf)
    total = len(reader.pages)
    if inicio < 1 or fin > total:
        raise ValueError(f"Rango inválido: el PDF tiene {total} páginas. Las páginas {inicio}-{fin} están fuera de rango.")
    if inicio > fin:
        raise ValueError("La página inicial no puede ser mayor que la final.")
    writer = PdfWriter()
    for i in range(inicio - 1, fin):
        writer.add_page(reader.pages[i])
    with open(salida, "wb") as f:
        writer.write(f)


